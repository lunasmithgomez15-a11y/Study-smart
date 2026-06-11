import json
import re
import random
import base64
import io
import os
import urllib.parse
import streamlit as st
from pypdf import PdfReader
from youtube_transcript_api import YouTubeTranscriptApi
from google import genai
from google.genai import types
from pydantic import BaseModel

# Safe fallbacks for optional document processing libraries
try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    import docx
    from docx import Document
except ImportError:
    docx = None
    Document = None

# --- LOCAL STORAGE FILE PATH ---
STORAGE_FILE = "storage.json"

def load_local_storage():
    """Loads saved quizzes and notes from the hard drive safely."""
    if os.path.exists(STORAGE_FILE):
        try:
            with open(STORAGE_FILE, "r", encoding="utf-8") as f:
                content = f.read().strip()
                if not content:
                    return {}
                return json.loads(content)
        except Exception:
            return {}
    return {}

def save_local_storage(data):
    """Saves binders permanently to the hard drive for offline access."""
    try:
        with open(STORAGE_FILE, "w", encoding="utf-8") as f:
            json.dump(data, f, ensure_ascii=False, indent=4)
    except Exception as e:
        st.error(f"Failed to write to local storage: {e}")

# --- PAGE CONFIGURATION ---
st.set_page_config(
    page_title="BrainCrunch Game Arena", 
    page_icon="🎮", 
    layout="centered"
)

# --- PYDANTIC BLUEPRINT MODELS ---
class QuizQuestion(BaseModel):
    question: str
    options: list[str]
    correct: str
    explanation: str

# --- INITIALIZE STATE VARIABLES ---
if "api_key" not in st.session_state:
    st.session_state.api_key = ""
if "questions" not in st.session_state:
    st.session_state.questions = []
if "review_notes" not in st.session_state:
    st.session_state.review_notes = ""
if "active_mode" not in st.session_state:
    st.session_state.active_mode = "Welcome" 
if "current_index" not in st.session_state:
    st.session_state.current_index = 0
if "score" not in st.session_state:
    st.session_state.score = 0
if "correct_answers_count" not in st.session_state:
    st.session_state.correct_answers_count = 0
if "answered" not in st.session_state:
    st.session_state.answered = False
if "selected_option" not in st.session_state:
    st.session_state.selected_option = None
if "streak" not in st.session_state:
    st.session_state.streak = 0
if "max_streak" not in st.session_state:
    st.session_state.max_streak = 0
if "leaderboard" not in st.session_state:
    st.session_state.leaderboard = []
if "quiz_host_persona" not in st.session_state:
    st.session_state.quiz_host_persona = "Enthusiastic School Teacher"
if "app_language" not in st.session_state:
    st.session_state.app_language = "English"

# Initialize binders directly from local hard drive storage file
if "nested_folders" not in st.session_state:
    st.session_state.nested_folders = load_local_storage()

# Student Lifeline Tracking Configurations
if "lifeline_5050_used" not in st.session_state:
    st.session_state.lifeline_5050_used = False
if "lifeline_hint_used" not in st.session_state:
    st.session_state.lifeline_hint_used = False
if "hidden_options" not in st.session_state:
    st.session_state.hidden_options = []

# AUTOMATIC SECRET API KEY CHECK 
if "gemini" in st.secrets:
    st.session_state.api_key = st.secrets["gemini"]["api_key"]

# --- SFX AUDIO INJECTION ---
def play_sfx(audio_url):
    st.components.v1.html(
        f"""
        <audio autoplay style="display:none;">
            <source src="{audio_url}" type="audio/mp3">
        </audio>
        """,
        height=0,
    )

CORRECT_SFX = "https://actions.google.com/sounds/v1/alarms/digital_watch_alarm_long.ogg"
WRONG_SFX = "https://actions.google.com/sounds/v1/cartoon/slide_whistle_down.ogg"

# --- FILE EXTRACTOR MECHANISM ---
def extract_text_from_file(file):
    filename = file.name.lower()
    text = ""
    if filename.endswith(".pdf"):
        reader = PdfReader(file)
        for page in reader.pages:
            text += page.extract_text() or ""
    elif filename.endswith(".pptx"):
        if Presentation:
            prs = Presentation(file)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
    elif filename.endswith(".docx"):
        if docx:
            doc = docx.Document(file)
            text += "\n".join([p.text for p in doc.paragraphs])
    elif filename.endswith(".txt"):
        text += file.read().decode("utf-8", errors="ignore")
    return text

def get_youtube_id(url):
    pattern = r'(?:https?:\/\/)?(?:www\.)?(?:youtube\.com\/(?:[^\/\n\s]+\/\S+\/|(?:v|e(?:mbed)?)\/|\S*?[?&]v=)|youtu\.be\/)([a-zA-Z0-9_-]{11})'
    match = re.search(pattern, url)
    return match.group(1) if match else None

def get_youtube_transcript(video_id):
    try:
        transcript_list = YouTubeTranscriptApi.get_transcript(video_id)
        return " ".join([item['text'] for item in transcript_list])
    except Exception:
        return None

# --- DOCX CONVERTER ENGINE ---
def build_docx_bytes(markdown_text):
    if not Document:
        return None
    doc = Document()
    doc.add_heading("BrainCrunch Study Reviewer Sheet", level=1)
    
    clean_lines = markdown_text.split("\n")
    for line in clean_lines:
        if line.startswith("## "):
            doc.add_heading(line.replace("## ", ""), level=2)
        elif line.startswith("### "):
            doc.add_heading(line.replace("### ", ""), level=3)
        elif line.startswith("- ") or line.startswith("* "):
            clean_bullet = line.replace("- ", "").replace("* ", "")
            doc.add_paragraph(clean_bullet, style="List Bullet")
        else:
            if line.strip():
                doc.add_paragraph(line)
                
    bio = io.BytesIO()
    doc.save(bio)
    return bio.getvalue()

# --- DEEP AI CORE ENGINES (PRO -> FLASH FALLBACK WITH OVERLOAD PROTECTION) ---
def generate_questions_with_ai(study_material, api_key, num_questions, persona, language):
    if not api_key:
        st.error("🛑 Offline/Configuration Limit: Generating new questions requires an internet connection and an API key.")
        return None

    lang_instruction = "All outputs must be written entirely in English."
    if language == "Tagalog / Filipino":
        lang_instruction = "CRITICAL: You are teaching a Filipino/Tagalog class. All questions, options, and explanations MUST be written in clear, natural Tagalog."

    prompt = f"""
    You are a smart game host playing with a student. Use this personality profile: "{persona}".
    {lang_instruction}
    
    TASK: Parse the provided text material thoroughly. Perform maximum deep information extraction—do not skip technical details, math variables, or niche terms. Create exactly {num_questions} high-quality multiple-choice questions.
    
    MATH RULE: If the question or options involve complex formulas, chemical equations, or mathematical expressions, format them cleanly using standard LaTeX notation (wrap with single '$' for inline math or double '$$' for large block equations) so they render beautifully on screen.
    
    Study Material:
    {study_material}
    """
    
    for target_model in ['gemini-2.5-flash', 'gemini-2.5-pro']:
        try:
            client = genai.Client(api_key=api_key)
            response = client.models.generate_content(
                model=target_model,
                contents=prompt,
                config=types.GenerateContentConfig(
                    response_mime_type="application/json",
                    response_schema=list[QuizQuestion],
                ),
            )
            return json.loads(response.text)
        except Exception as e:
            if any(err in str(e).upper() for err in ["503", "429", "UNAVAILABLE", "RESOURCE_EXHAUSTED", "LIMIT"]):
                continue
            st.error(f"🛑 AI System Error: {e}")
            return None
    
    st.error("⏳ All available AI models are currently busy or rate-limited. Please wait a moment and try clicking process again.")
    return None

def generate_summary_with_ai(study_material, api_key, persona, language):
    if not api_key:
        st.error("🛑 Offline/Configuration Limit: Generating new notes requires an internet connection and an API key.")
        return None

    lang_instruction = "Write the summary sheet in English."
    if language == "Tagalog / Filipino":
        lang_instruction = "CRITICAL: Write the entire summary sheet in fluent Tagalog/Filipino language."

    prompt = f"""
    You are an expert academic tutor with this personality profile: "{persona}".
    {lang_instruction}
    
    TASK: Read the uploaded file text and perform an exhaustive deep information extraction. Pull out all definitions, key historical dates, core concepts, formulas, and laws. 
    
    FORMAT RULES:
    - Organize using neat Markdown bullet points and bold headers.
    - For mathematical formulas, equations, matrices, or variable fractions, use beautiful, readable standard LaTeX tags ($...$ or $$...$$). Make numbers and step-by-step math breakdowns perfectly organized.
    
    Study Material:
    {study_material}
    """
    
    for target_model in ['gemini-2.5-flash', 'gemini-2.5-pro']:
        try:
            client = genai.Client(api_key=api_key)
            response = client.models.generate_content(
                model=target_model,
                contents=prompt
            )
            return response.text
        except Exception as e:
            if any(err in str(e).upper() for err in ["503", "429", "UNAVAILABLE", "RESOURCE_EXHAUSTED", "LIMIT"]):
                continue
            st.error(f"🛑 AI Summary Error: {e}")
            return None
            
    st.error("⏳ All available AI models are currently busy or rate-limited. Please wait a moment and try clicking process again.")
    return None

# --- APP FRONTEND ---
st.markdown("<h1 style='text-align: center;'>🧠 BrainCrunch Studio Pro</h1>", unsafe_allow_html=True)

# --- SIDEBAR CONTROL UNIT ---
with st.sidebar:
    st.header("⚙️ App Mode")
    user_role = st.selectbox("I am a...", ["Player / Student", "Admin / Creator"])
    
    if user_role == "Admin / Creator":
        st.write("---")
        st.subheader("🔑 Admin Access")
        admin_pass = st.text_input("Enter Admin Password:", type="password")
        if admin_pass == "studio123":
            st.success("Access Verified! 🛠️")
            
            st.markdown("### 🤫 Secret Admin Feature Hub")
            st.session_state.api_key = st.text_input("System Gemini API Key:", value=st.session_state.api_key, type="password")
            
            st.subheader("🧙‍♂️ Game Host Persona")
            st.session_state.quiz_host_persona = st.selectbox(
                "Choose AI Game Master:",
                ["Enthusiastic School Teacher", "Sarcastic Pirate Coach", "Strict Drill Sergeant", "Whimsical Fantasy Wizard"]
            )
        elif admin_pass:
            st.error("Incorrect Password!")

    st.write("---")
    st.header("🎮 Student Core Features")
        
    st.subheader("🌐 Language Filter")
    st.session_state.app_language = st.selectbox("Select Study Language:", ["English", "Tagalog / Filipino"])
        
    st.write("---")
    st.subheader("🧙‍♂️ AI Content Engine (Requires Internet)")
    output_type = st.radio("What should the AI build?", ["Gamified Quiz Decks", "Clean Review Summaries"])
    
    if output_type == "Gamified Quiz Decks":
        question_count = st.number_input("How many questions?", min_value=1, max_value=150, value=5, step=5)
    
    st.write("---")
    input_mode = st.radio("Input Source:", ["Upload Files (PDF, PPTX, DOCX, TXT)", "Voice Lesson Record / Audio Note", "YouTube Video Link"])
    
    study_text = ""
    triggered_generation = False
    
    if input_mode == "Upload Files (PDF, PPTX, DOCX, TXT)":
        uploaded_files = st.file_uploader("Drop slides or files:", type=["pdf", "pptx", "docx", "txt"], accept_multiple_files=True)
        if uploaded_files and st.button("🚀 Process Study Material", use_container_width=True):
            with st.spinner("Extracting content strings... 📂"):
                for f in uploaded_files:
                    study_text += extract_text_from_file(f) + "\n"
                triggered_generation = True

    elif input_mode == "Voice Lesson Record / Audio Note":
        recorded_audio = st.file_uploader("Upload audio lesson clip:", type=["mp3", "wav", "m4a", "ogg"])
        if recorded_audio and st.button("🎙️ Process Lesson Audio Track", use_container_width=True):
            if not st.session_state.api_key:
                st.error("🛑 API Key missing. Please make sure the Admin has provided a valid Gemini API Key configuration setup.")
            else:
                with st.spinner("Transcribing lesson audio... 💬"):
                    try:
                        client = genai.Client(api_key=st.session_state.api_key)
                        audio_upload_res = client.files.upload(file=recorded_audio, mime_type=recorded_audio.type)
                        tx_prompt = "Transcribe the following lecture audio track exactly, keeping all numbers and core concepts sharp."
                        
                        # Added structured model fallbacks for audio transcription process
                        tx_response = None
                        for model_option in ["gemini-2.5-flash", "gemini-2.5-pro"]:
                            try:
                                tx_response = client.models.generate_content(model=model_option, contents=[audio_upload_res, tx_prompt])
                                break
                            except Exception as sub_err:
                                if any(err in str(sub_err).upper() for err in ["503", "429", "UNAVAILABLE", "RESOURCE_EXHAUSTED"]):
                                    continue
                                raise sub_err
                        
                        if tx_response:
                            study_text = tx_response.text
                            triggered_generation = True
                        else:
                            st.error("⏳ Audio engine is currently overwhelmed. Please wait a moment and try again.")
                    except Exception as audio_err:
                        st.error(f"Audio processing failure checklist: {audio_err}")

    elif input_mode == "YouTube Video Link":
        yt_url = st.text_input("YouTube Video URL:")
        if yt_url and st.button("🎬 Process Video Streams", use_container_width=True):
            with st.spinner("Analyzing transcript... 🍿"):
                v_id = get_youtube_id(yt_url)
                if v_id:
                    extracted_yt = get_youtube_transcript(v_id)
                    if extracted_yt:
                        study_text = extracted_yt
                        triggered_generation = True
    
    if triggered_generation and study_text:
        if output_type == "Gamified Quiz Decks":
            ai_qs = generate_questions_with_ai(study_text, st.session_state.api_key, question_count, st.session_state.quiz_host_persona, st.session_state.app_language)
            if ai_qs:
                st.session_state.questions = ai_qs
                st.session_state.active_mode = "Quiz"
                st.session_state.current_index = 0
                st.session_state.score = 0
                st.session_state.correct_answers_count = 0
                st.session_state.streak = 0
                st.session_state.max_streak = 0
                st.session_state.answered = False
                st.session_state.lifeline_5050_used = False
                st.session_state.lifeline_hint_used = False
                st.session_state.hidden_options = []
                st.rerun()
        else:
            ai_notes = generate_summary_with_ai(study_text, st.session_state.api_key, st.session_state.quiz_host_persona, st.session_state.app_language)
            if ai_notes:
                st.session_state.review_notes = ai_notes
                st.session_state.active_mode = "Reviewer"
                st.rerun()

    st.write("---")
    st.subheader("📁 Save to Binders (Works Offline)")
    path_input = st.text_input("Folder Path:", value="Q1 / Science / Biology")
    if st.button("📂 File Current Data Into Path", use_container_width=True):
        if st.session_state.active_mode == "Quiz" and st.session_state.questions:
            st.session_state.nested_folders[path_input] = {"type": "Quiz", "data": list(st.session_state.questions)}
            save_local_storage(st.session_state.nested_folders)
            st.toast(f"Quiz saved to hard drive under: {path_input}!")
            st.rerun()
        elif st.session_state.active_mode == "Reviewer" and st.session_state.review_notes:
            st.session_state.nested_folders[path_input] = {"type": "Reviewer", "data": st.session_state.review_notes}
            save_local_storage(st.session_state.nested_folders)
            st.toast(f"Summary Reviewer saved to hard drive under: {path_input}!")
            st.rerun()

    # STUDENT BINDER DROPDOWN
    st.write("---")
    st.header("🗂️ Study Binders (Works Offline)")
    if st.session_state.nested_folders:
        selected_path = st.selectbox("Choose a study track to open:", list(st.session_state.nested_folders.keys()))
        if st.button("🎮 Load Selected Content", use_container_width=True):
            saved_item = st.session_state.nested_folders[selected_path]
            if saved_item["type"] == "Quiz":
                st.session_state.questions = list(saved_item["data"])
                st.session_state.active_mode = "Quiz"
                st.session_state.current_index = 0
                st.session_state.score = 0
                st.session_state.correct_answers_count = 0
                st.session_state.streak = 0
                st.session_state.max_streak = 0
                st.session_state.answered = False
                st.session_state.lifeline_5050_used = False
                st.session_state.lifeline_hint_used = False
                st.session_state.hidden_options = []
            else:
                st.session_state.review_notes = saved_item["data"]
                st.session_state.active_mode = "Reviewer"
            st.rerun()
    else:
        st.caption("No custom subject tracks saved yet.")

# --- MAIN RUNTIME ARENA ---

if st.session_state.active_mode == "Welcome":
    st.info("👋 **Welcome to the Game Arena!**\n\nUse the sidebar panel to generate flash quizzes, read summary sheets, or load saved topics out of your Study Binders folder cabinet tracker!")

elif st.session_state.active_mode == "Reviewer":
    st.markdown("## 📑 Smart Summary Reviewer Sheet")
    st.write("---")
    st.markdown(st.session_state.review_notes)
    st.write("---")
    
    # EXPORT HUB INTERFACE
    st.markdown("### 📤 Student Export Hub")
    st.caption("Copy this summary sheet directly to your computer clipboard or extract it into your word processing documents!")
    
    col_exp1, col_exp2 = st.columns(2)
    with col_exp1:
        st.text_area("📋 Copy to Clipboard Panel (Select All -> Copy):", value=st.session_state.review_notes, height=80)
        
    with col_exp2:
        if Document:
            docx_data = build_docx_bytes(st.session_state.review_notes)
            st.download_button(
                label="📄 Export to Word File (.docx)",
                data=docx_data,
                file_name="BrainCrunch_Study_Sheet.docx",
                mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                use_container_width=True
            )
        else:
            st.download_button(
                label="📝 Export to Plain Text (.txt)",
                data=st.session_state.review_notes,
                file_name="BrainCrunch_Study_Sheet.txt",
                mime="text/plain",
                use_container_width=True
            )

    st.write("---")
    
    # EXTERNAL DIRECTIONAL HUBS
    st.markdown("### 🔍 Need a Video Explainer for This Lesson? (Requires Internet)")
    try:
        clean_topic_query = urllib.parse.quote(selected_path.replace("/", " "))
    except Exception:
        clean_topic_query = "study+lesson"
        
    yt_search_url = f"https://www.youtube.com/results?search_query={clean_topic_query}+lesson+explanation"
    google_search_url = f"https://www.google.com/search?q={clean_topic_query}+educational+guide"
    
    col_nav1, col_nav2 = st.columns(2)
    with col_nav1:
        st.link_button("📺 Search Lessons on YouTube", yt_search_url, use_container_width=True)
    with col_nav2:
        st.link_button("🌐 Search Articles on Google", google_search_url, use_container_width=True)
        
    st.write("---")
    if st.button("🏠 Back to Home Screen", key="back_home_rev", use_container_width=True):
        st.session_state.active_mode = "Welcome"
        st.rerun()

elif st.session_state.active_mode == "Quiz":
    idx = st.session_state.current_index
    if idx < len(st.session_state.questions):
        current_q = st.session_state.questions[idx]
        
        col1, col2, col3 = st.columns(3)
        with col1:
            st.markdown(f"<h4 style='text-align:center;color:#ffaa00;'>🏆 Score</h4><p style='text-align:center;font-size:24px;font-weight:bold;'>{st.session_state.score}</p>", unsafe_allow_html=True)
        with col2:
            st.markdown(f"<h4 style='text-align:center;color:#ff5555;'>🔥 Streak</h4><p style='text-align:center;font-size:24px;font-weight:bold;'>{st.session_state.streak}</p>", unsafe_allow_html=True)
        with col3:
            st.markdown(f"<h4 style='text-align:center;color:#55afb5;'>🎯 Stage</h4><p style='text-align:center;font-size:24px;font-weight:bold;'>{idx + 1} / {len(st.session_state.questions)}</p>", unsafe_allow_html=True)
            
        st.progress((idx) / len(st.session_state.questions))
        st.write("---")
        
        # STUDENT LIFELINES HUB
        if not st.session_state.answered:
            st.markdown("##### 🆘 Student Lifelines")
            col_life1, col_life2 = st.columns(2)
            with col_life1:
                if st.button("⚖️ Use 50/50", key=f"life_50_{idx}", disabled=st.session_state.lifeline_5050_used, use_container_width=True):
                    st.session_state.lifeline_5050_used = True
                    wrong_choices = [opt for opt in current_q['options'] if opt[0] != current_q['correct']]
                    if wrong_choices:
                        st.session_state.hidden_options = [random.choice(wrong_choices)]
                    st.rerun()
            with col_life2:
                if st.button("💡 Ask AI Guide Clue", key=f"life_hint_{idx}", disabled=st.session_state.lifeline_hint_used, use_container_width=True):
                    st.session_state.lifeline_hint_used = True
                    st.session_state.show_hint_text = True
            if "show_hint_text" in st.session_state and st.session_state.show_hint_text:
                st.caption(f"🤖 *AI Companion Clue:* Look carefully at matching rules and definitions!")
            st.write("---")
        
        st.markdown(f"### ❓ {current_q['question']}")
        
        for option in current_q['options']:
            if option in st.session_state.hidden_options:
                continue 
            if not st.session_state.answered:
                if st.button(option, key=f"btn_{idx}_{option}", use_container_width=True):
                    st.session_state.answered = True
                    st.session_state.selected_option = option
                    st.rerun()
                    
        if st.session_state.answered:
            user_letter = st.session_state.selected_option[0]
            correct_letter = current_q["correct"]
            
            if user_letter == correct_letter:
                st.markdown(f"<div style='background-color:#d4edda; color:#155724; border-radius:10px; padding:15px; border-left:6px solid #28a745; margin-bottom:15px;'><h4 style='margin:0;'>🎉 EXCELLENT HIT!</h4><p style='margin:5px 0 0 0;'>You chose: <b>{st.session_state.selected_option}</b></p></div>", unsafe_allow_html=True)
                if f"scored_{idx}" not in st.session_state:
                    st.session_state.score += 10 + (st.session_state.streak * 2)
                    st.session_state.streak += 1
                    st.session_state.correct_answers_count += 1
                    if st.session_state.streak > st.session_state.max_streak:
                        st.session_state.max_streak = st.session_state.streak
                    st.session_state[f"scored_{idx}"] = True
            else:
                st.markdown(f"<div style='background-color:#f8d7da; color:#721c24; border-radius:10px; padding:15px; border-left:6px solid #dc3545; margin-bottom:15px;'><h4 style='margin:0;'>💔 DEFLECTED</h4><p style='margin:5px 0 0 0;'>You chose <b>{user_letter}</b>. Correct path: <b>{correct_letter}</b>.</p></div>", unsafe_allow_html=True)
                st.session_state.streak = 0
                
            st.markdown(f"**🧠 Memory Scoop ({st.session_state.quiz_host_persona}):**")
            st.write(current_q['explanation'])
            
            # DYNAMIC REDIRECTION HUD
            st.write("")
            st.markdown("#### 📺 Confused About This Question? (Requires Internet)")
            query_keywords = re.sub(r'[^\w\s]', '', current_q['question'])
            encoded_query = urllib.parse.quote(query_keywords)
            
            question_yt_url = f"https://www.youtube.com/results?search_query={encoded_query}"
            question_google_url = f"https://www.google.com/search?q={encoded_query}"
            
            col_qnav1, col_qnav2 = st.columns(2)
            with col_qnav1:
                st.link_button("📺 Watch Video Explainers", question_yt_url, use_container_width=True)
            with col_qnav2:
                st.link_button("🌐 Google Text Breakdowns", question_google_url, use_container_width=True)
            
            st.write("---")
            if st.button("➡️ Advance to Next Level", use_container_width=True):
                st.session_state.current_index += 1
                st.session_state.answered = False
                st.session_state.selected_option = None
                st.session_state.hidden_options = []
                if "show_hint_text" in st.session_state:
                    del st.session_state.show_hint_text
                st.rerun()
    else:
        st.success("🏆 **CAMP RUN COMPLETED!** 🏆")
        st.subheader("📊 Performance Scorecard")
        
        total_qs = len(st.session_state.questions)
        accuracy_pct = int((st.session_state.correct_answers_count / total_qs) * 100) if total_qs > 0 else 0
        
        col_an1, col_an2, col_an3 = st.columns(3)
        with col_an1:
            st.metric(label="🎖️ Score", value=f"{st.session_state.score} pts")
        with col_an2:
            st.metric(label="🎯 Accuracy", value=f"{accuracy_pct}%")
        with col_an3:
            st.metric(label="🔥 Max Streak", value=f"{st.session_state.max_streak}")
            
        player_name = st.text_input("Enter name for scoreboard matrix:", value="Player 1")
        if st.button("💾 Log Score", use_container_width=True):
            st.session_state.leaderboard.append({"name": player_name, "score": st.session_state.score, "accuracy": accuracy_pct})
            st.toast("Score logged into local system matrix!")
            st.rerun()
            
        if st.session_state.leaderboard:
            st.write("### 🏁 Current Rankings")
            sorted_board = sorted(st.session_state.leaderboard, key=lambda x: x['score'], reverse=True)
            for place, entry in enumerate(sorted_board, 1):
                st.write(f"**#{place}** {entry['name']} — `{entry['score']} pts` ({entry.get('accuracy', 0)}% Accuracy)")
        
        if st.button("🔄 Replay Session", use_container_width=True):
            st.session_state.current_index = 0
            st.session_state.score = 0
            st.session_state.correct_answers_count = 0
            st.session_state.streak = 0
            st.session_state.max_streak = 0
            st.session_state.answered = False
            st.session_state.selected_option = None
            st.session_state.lifeline_5050_used = False
            st.session_state.lifeline_hint_used = False
            st.session_state.hidden_options = []
            st.rerun()
